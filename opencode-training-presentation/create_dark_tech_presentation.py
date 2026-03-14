#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
sys.path.insert(0, '/home/jabing/.local/lib/python3.14/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 暗黑科技风格（Dark Tech Style）
BACKGROUND_COLOR = RGBColor(10, 10, 10)      # 深黑色
CARD_COLOR = RGBColor(30, 30, 35)          # 深灰卡片
PRIMARY_COLOR = RGBColor(0, 212, 255)      # 霓虹蓝
ACCENT_COLOR = RGBColor(0, 255, 255)       # 霓虹青
HIGHLIGHT_COLOR = RGBColor(255, 0, 255)     # 霓虹紫
TEXT_COLOR = RGBColor(255, 255, 255)        # 白色
CODE_COLOR = RGBColor(0, 255, 136)         # 亮绿代码

# 字体
TITLE_FONT = "Arial Black"
BODY_FONT = "Arial"
CODE_FONT = "Consolas"

def add_dark_background(slide):
    """添加暗黑背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR

def add_title_slide(prs, title, subtitle, tag):
    """暗黑科技风标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_dark_background(slide)
    
    # 装饰线条（科技感）
    for i in range(5):
        line = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3 + i * 0.15), Inches(0.5 + i * 0.3),
            Inches(9.4 - i * 0.1), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.fill.fore_color.brightness = 0.5 + i * 0.1
        line.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(76)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = HIGHLIGHT_COLOR
    subtitle_para.font.name = BODY_FONT
    
    # 标签（发光效果）
    tag_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.5), Inches(4), Inches(0.6)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = CARD_COLOR
    tag_box.line.color.rgb = ACCENT_COLOR
    tag_box.line.width = Pt(2)
    
    tag_text = tag_box.text_frame
    tag_text.text = tag
    tag_para = tag_text.paragraphs[0]
    tag_para.alignment = PP_ALIGN.CENTER
    tag_para.font.size = Pt(24)
    tag_para.font.bold = True
    tag_para.font.color.rgb = ACCENT_COLOR
    tag_para.font.name = TITLE_FONT

def add_content_card_slide(prs, title, content_items):
    """内容卡片页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_dark_background(slide)
    
    # 装饰元素
    corner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.5), Inches(4.5), Inches(1.5), Inches(1.5)
    )
    corner.fill.solid()
    corner.fill.fore_color.rgb = PRIMARY_COLOR
    corner.fill.fore_color.brightness = 0.2
    corner.line.fill.background()
    
    # 标题栏
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = CARD_COLOR
    title_bar.line.color.rgb = PRIMARY_COLOR
    title_bar.line.width = Pt(2)
    
    title_text = title_bar.text_frame
    title_text.text = title
    title_para = title_text.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 内容卡片
    y_pos = 1.3
    for item in content_items:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = ACCENT_COLOR
        card.line.width = Pt(1)
        
        text = card.text_frame
        text.text = item
        para = text.paragraphs[0]
        para.font.size = Pt(22)
        para.font.color.rgb = TEXT_COLOR
        para.font.name = BODY_FONT
        
        y_pos += 0.9

def add_code_slide(prs, title, code_blocks):
    """代码展示页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_dark_background(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 代码块
    y_pos = 1.2
    for code_block in code_blocks:
        # 代码框
        code_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.6)
        )
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(20, 20, 20)
        code_box.line.color.rgb = CODE_COLOR
        code_box.line.width = Pt(2)
        
        # 代码文本
        code_text = code_box.text_frame
        code_text.text = code_block
        code_para = code_text.paragraphs[0]
        code_para.font.size = Pt(16)
        code_para.font.color.rgb = CODE_COLOR
        code_para.font.name = CODE_FONT
        
        y_pos += 0.75

def add_feature_grid_slide(prs, title, features):
    """特性网格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_dark_background(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 2x3 特性网格
    positions = [
        (Inches(0.5), Inches(1.2)),
        (Inches(5), Inches(1.2)),
        (Inches(0.5), Inches(2.4)),
        (Inches(5), Inches(2.4)),
        (Inches(0.5), Inches(3.6)),
        (Inches(5), Inches(3.6))
    ]
    
    for feature, pos in zip(features, positions):
        x, y = pos
        
        # 特性卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = ACCENT_COLOR
        card.line.width = Pt(1)
        
        # 文本
        text = card.text_frame
        text.text = feature
        para = text.paragraphs[0]
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = TEXT_COLOR
        para.font.name = BODY_FONT

def add_flow_slide(prs, title, steps):
    """流程图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_dark_background(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 流程节点
    x_pos = 0.5
    for i, step in enumerate(steps):
        # 节点
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.3), Inches(1.8), Inches(2.2)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = CARD_COLOR
        node.line.color.rgb = HIGHLIGHT_COLOR
        node.line.width = Pt(2)
        
        # 文本
        text = node.text_frame
        text.word_wrap = True
        text.text = step
        para = text.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = TEXT_COLOR
        para.font.name = BODY_FONT
        
        x_pos += 1.9
        
        # 连接线
        if i < len(steps) - 1:
            conn = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos - 0.3), Inches(2.5),
                Inches(0.5), Inches(0.1)
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = ACCENT_COLOR

def add_qa_slide(prs, title, subtitle):
    """问答页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_dark_background(slide)
    
    # 装饰圆
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.5), Inches(2.5), Inches(2.5), Inches(2.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = HIGHLIGHT_COLOR
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(88)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = ACCENT_COLOR
    subtitle_para.font.name = BODY_FONT

def create_dark_tech_presentation():
    """创建暗黑科技风演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 标题页
    add_title_slide(
        prs,
        "OpenCode 安装\n与配置指南",
        "从零开始打造 AI 编码助手",
        "< 部门内部培训 >"
    )
    
    # 2. 目录
    add_content_card_slide(
        prs,
        "课程大纲",
        [
            "1. Node.js 安装",
            "2. OpenCode 安装",
            "3. Easy-OpenCode 部署",
            "4. DeepSeek API Key 配置",
            "5. 开始使用 Easy-OpenCode",
            "6. 常见问题与解决方案"
        ]
    )
    
    # 3. 第一章：Node.js
    add_code_slide(
        prs,
        "第一章：Node.js 安装",
        [
            "# 官方安装包（推荐）",
            "访问: https://nodejs.org/",
            "下载 LTS 版本",
            "Windows: 双击 .msi 文件",
            "macOS: 双击 .pkg 文件",
            "Linux: 按照官网指南操作",
            "",
            "# 验证安装",
            "node --version",
            "npm --version"
        ]
    )
    
    # 4. 第二章：OpenCode
    add_code_slide(
        prs,
        "第二章：OpenCode 安装",
        [
            "# 全局安装（推荐）",
            "npm install -g opencode",
            "",
            "# 验证安装",
            "opencode --version",
            "opencode --help"
        ]
    )
    
    # 5. 第三章：EOC 特性
    add_feature_grid_slide(
        prs,
        "第三章：Easy-OpenCode 核心特性",
        [
            "14 个专业 AI Agent",
            "50+ 工作流技能",
            "33 个快捷命令",
            "自动化钩子工作流",
            "TDD 驱动开发",
            "内置安全审查"
        ]
    )
    
    # 6. EOC 安装
    add_code_slide(
        prs,
        "Easy-OpenCode 安装",
        [
            "# 全局安装（推荐）",
            "npm install -g easy-opencode",
            "",
            "# 在项目中安装",
            "eoc-install",
            "或",
            "npx easy-opencode install"
        ]
    )
    
    # 7. API Key 配置
    add_code_slide(
        prs,
        "第四章：DeepSeek API Key 配置",
        [
            "# 环境变量配置（推荐）",
            "export DEEPSEEK_API_KEY=\"your-api-key-here\"",
            "source ~/.bashrc  # 或 source ~/.zshrc",
            "",
            "# OpenCode 配置文件（持久化）",
            "# 编辑: ~/.config/opencode/opencode.json",
            "# 添加:",
            '{\"env\": {\"DEEPSEEK_API_KEY\": \"your-api-key-here\"}}',
            "",
            "# 安全注意事项",
            "# 将 .env 添加到 .gitignore",
            "# 定期轮换 API Key"
        ]
    )
    
    # 8. EOC 命令
    add_content_card_slide(
        prs,
        "第五章：常用 EOC 命令",
        [
            "/plan - 创建实现计划",
            "/tdd - TDD 开发流程",
            "/code-review - 代码审查",
            "/security - 安全审查",
            "/build-fix - 修复构建错误",
            "/architect - 架构设计",
            "/e2e - 生成 E2E 测试",
            "/refactor-clean - 清理死代码"
        ]
    )
    
    # 9. 实战演示
    add_flow_slide(
        prs,
        "实战演示：30 分钟完整开发流程",
        [
            "步骤 1: 规划\n/plan\n5分钟",
            "步骤 2: TDD 开发\n/tdd\n15分钟",
            "步骤 3: 安全审查\n/security\n3分钟",
            "步骤 4: 代码审查\n/code-review\n3分钟",
            "步骤 5: E2E 测试\n/e2e\n5分钟"
        ]
    )
    
    # 10. 常见问题
    add_content_card_slide(
        prs,
        "第六章：常见问题与解决方案",
        [
            "Q: Node.js 安装失败",
            "A: 检查 PATH 环境变量，重新运行安装程序",
            "Q: OpenCode 安装失败",
            "A: 检查网络连接，清除 npm 缓存后重试",
            "Q: API Key 不生效",
            "A: 验证环境变量，重启 OpenCode",
            "Q: Agent 不可用",
            "A: 确认 EOC 安装成功，重新安装"
        ]
    )
    
    # 11. 总结
    add_feature_grid_slide(
        prs,
        "恭喜！您已掌握：",
        [
            "✅ Node.js 安装配置",
            "✅ OpenCode 安装验证",
            "✅ EOC 部署配置",
            "✅ API Key 安全设置",
            "✅ EOC 命令熟练使用",
            "✅ 工作流程完全掌握"
        ]
    )
    
    # 12. 下一步
    add_content_card_slide(
        prs,
        "下一步行动",
        [
            "🚀 立即实践：创建测试项目，尝试第一个命令",
            "📚 深入学习：阅读 AGENTS.md，探索 skills/ 目录",
            "💬 团队协作：分享经验，建立最佳实践",
            "🔄 持续改进：定期技术交流，追踪效率提升"
        ]
    )
    
    # 13. Q&A
    add_qa_slide(
        prs,
        "Q & A",
        "感谢聆听！如有疑问，欢迎提问 💡"
    )
    
    return prs

if __name__ == "__main__":
    prs = create_dark_tech_presentation()
    output_path = "opencode-training-presentation/OpenCode安装与配置指南-暗黑科技版.pptx"
    prs.save(output_path)
    print(f"✅ 暗黑科技风演示文稿已生成: {output_path}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    print(f"🎨 设计风格: Dark Tech (暗黑科技)")
    print(f"🎯 配色方案:")
    print(f"   - 背景: #0A0A0A (深黑)")
    print(f"   - 主色: #00D4FF (霓虹蓝)")
    print(f"   - 辅助: #00FFFF (霓虹青)")
    print(f"   - 强调: #FF00FF (霓虹紫)")
    print(f"   - 代码: #00FF88 (亮绿)")
    print(f"✨ 特性: 暗黑主题、霓虹强调、代码高亮、极简布局")
