#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
sys.path.insert(0, '/home/jabing/.local/lib/python3.14/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Coral Energy 主题
PRIMARY_COLOR = RGBColor(249, 97, 103)      # Coral
SECONDARY_COLOR = RGBColor(249, 231, 149)   # Gold
BACKGROUND_COLOR = RGBColor(47, 60, 126)    # Navy
TEXT_COLOR = RGBColor(255, 255, 255)        # White
CARD_COLOR = RGBColor(255, 255, 255)        # White

TITLE_FONT = "Arial Black"
BODY_FONT = "Arial"

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title_slide(prs, title, subtitle, tag):
    """标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_gradient_background(slide, PRIMARY_COLOR, BACKGROUND_COLOR)
    
    # 装饰圆
    for i in range(3):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(0.5 + i * 0.8), 
            Inches(4 + i * 0.3), 
            Inches(0.8 + i * 0.3), 
            Inches(0.8 + i * 0.3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = SECONDARY_COLOR
        circle.fill.fore_color.brightness = 0.4
        circle.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(88)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.font.name = BODY_FONT
    
    # 标签
    tag_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.3), Inches(3), Inches(0.5)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = SECONDARY_COLOR
    tag_box.line.fill.background()
    
    tag_text = tag_box.text_frame
    tag_text.text = tag
    tag_para = tag_text.paragraphs[0]
    tag_para.alignment = PP_ALIGN.CENTER
    tag_para.font.size = Pt(22)
    tag_para.font.bold = True
    tag_para.font.color.rgb = BACKGROUND_COLOR
    tag_para.font.name = BODY_FONT

def add_card_grid_slide(prs, title, items, icon_prefix=""):
    """卡片网格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(52)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 2x3 卡片网格
    positions = [
        (Inches(0.5), Inches(1.5)),
        (Inches(5), Inches(1.5)),
        (Inches(0.5), Inches(3)),
        (Inches(5), Inches(3)),
        (Inches(0.5), Inches(4.5)),
        (Inches(5), Inches(4.5))
    ]
    
    for item, pos in zip(items, positions):
        x, y = pos
        
        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = PRIMARY_COLOR
        card.line.width = Pt(3)
        
        # 编号圆
        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.3), Inches(0.6), Inches(0.6)
        )
        num.fill.solid()
        num.fill.fore_color.rgb = PRIMARY_COLOR
        num.line.fill.background()
        
        # 文本
        text = card.text_frame
        text.text = f"  {item}"
        para = text.paragraphs[0]
        para.font.size = Pt(22)
        para.font.bold = True
        para.font.color.rgb = BACKGROUND_COLOR
        para.font.name = BODY_FONT

def add_two_column_slide(prs, title, left_items, right_items):
    """双栏页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.8)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = PRIMARY_COLOR
    title_box.line.fill.background()
    
    title_text = title_box.text_frame
    title_text.text = title
    title_para = title_text.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.font.name = TITLE_FONT
    
    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.4), Inches(3.8))
    left_frame = left_box.text_frame
    
    for i, line in enumerate(left_items):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = BACKGROUND_COLOR
        p.font.name = BODY_FONT
        p.space_before = Pt(6)
        p.space_after = Pt(4)
    
    # 分隔条
    divider = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.5), Inches(0.1), Inches(3.8)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = SECONDARY_COLOR
    
    # 右栏
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(3.8))
    right_frame = right_box.text_frame
    
    for i, line in enumerate(right_items):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = BACKGROUND_COLOR
        p.font.name = BODY_FONT
        p.space_before = Pt(6)
        p.space_after = Pt(4)

def add_feature_boxes_slide(prs, title, features):
    """特性盒子页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = SECONDARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 特性盒子
    for i, feature in enumerate(features):
        row = i // 2
        col = i % 2
        
        x = Inches(0.5) if col == 0 else Inches(5)
        y = Inches(1.3) + Inches(row * 1.4)
        
        # 盒子
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_COLOR
        box.line.color.rgb = SECONDARY_COLOR
        box.line.width = Pt(2)
        
        # 文本
        text = box.text_frame
        text.text = feature
        para = text.paragraphs[0]
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = BACKGROUND_COLOR
        para.font.name = BODY_FONT

def add_command_slide(prs, title, commands):
    """命令页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, BACKGROUND_COLOR, RGBColor(30, 40, 90))
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = SECONDARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 命令条
    y_pos = 1.2
    for i, cmd in enumerate(commands):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.65)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = CARD_COLOR
        bar.line.color.rgb = SECONDARY_COLOR
        bar.line.width = Pt(2)
        
        # 文本
        text = bar.text_frame
        text.text = f"  {cmd}"
        para = text.paragraphs[0]
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = BACKGROUND_COLOR
        para.font.name = BODY_FONT
        
        y_pos += 0.8

def add_flow_slide(prs, title, steps):
    """流程页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 流程节点
    x_pos = 0.5
    for i, step in enumerate(steps):
        # 节点
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.5), Inches(1.8), Inches(2.8)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = PRIMARY_COLOR
        node.line.color.rgb = SECONDARY_COLOR
        node.line.width = Pt(3)
        
        # 文本
        text = node.text_frame
        text.word_wrap = True
        text.text = step
        para = text.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = TEXT_COLOR
        para.font.name = BODY_FONT
        
        x_pos += 1.9
        
        # 箭头（除了最后一个）
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(x_pos - 0.35), 
                Inches(2.8), 
                Inches(0.7), 
                Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SECONDARY_COLOR

def add_qa_slide(prs, title, subtitle):
    """问答页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, PRIMARY_COLOR, BACKGROUND_COLOR)
    
    # 装饰
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.5), Inches(2.5), Inches(2.5), Inches(2.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY_COLOR
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(96)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.font.name = BODY_FONT

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 标题页
    add_title_slide(
        prs,
        "OpenCode 安装\n与配置指南",
        "从零开始打造 AI 编码助手",
        "🎯 部门内部培训"
    )
    
    # 2. 目录
    add_card_grid_slide(
        prs,
        "课程大纲",
        [
            "Node.js 安装配置",
            "OpenCode 安装",
            "EOC 部署",
            "API Key 配置",
            "EOC 使用",
            "常见问题"
        ]
    )
    
    # 3. 第一章：Node.js
    add_two_column_slide(
        prs,
        "第一章：Node.js 安装",
        [
            "什么是 Node.js？",
            "• Chrome V8 引擎运行时",
            "• 跨平台支持",
            "• npm 包管理器",
            "• OpenCode 基础",
            "",
            "官方安装（推荐）",
            "• 访问 nodejs.org",
            "• 下载 LTS 版本",
            "• 运行安装程序"
        ],
        [
            "Linux 命令行安装",
            "• Ubuntu/Debian:",
            "  curl -fsSL...",
            "• sudo apt-get install",
            "",
            "验证安装",
            "• node --version",
            "• npm --version"
        ]
    )
    
    # 4. 第二章：OpenCode
    add_two_column_slide(
        prs,
        "第二章：OpenCode 安装",
        [
            "OpenCode 简介",
            "• 强大的 AI 编码助手",
            "• 13 个专业 AI Agent",
            "• 50+ 工作流技能",
            "• 33 个命令",
            "",
            "Homebrew 安装",
            "• brew install opencode",
            "• 适用于 macOS/Linux",
            "• 自动配置"
        ],
        [
            "NPM 全局安装",
            "• npm install -g opencode",
            "• 适用于所有平台",
            "• 验证安装",
            "",
            "验证成功标志",
            "• 显示帮助信息",
            "• Agent 列表可见",
            "• 命令列表完整"
        ]
    )
    
    # 5. 第三章：EOC 介绍
    add_feature_boxes_slide(
        prs,
        "第三章：Easy-OpenCode 特性",
        [
            "🤖 14 个专业 AI Agent",
            "⚡ 50+ 工作流技能",
            "🛠️ 33 个快捷命令",
            "🔗 自动化钩子工作流",
            "📦 一键安装开箱即用",
            "🎯 TDD 驱动确保质量"
        ]
    )
    
    # 6. EOC 安装
    add_two_column_slide(
        prs,
        "Easy-OpenCode 安装方式",
        [
            "全局安装（推荐）⭐",
            "• npm install -g easy-opencode",
            "• eoc-install --version",
            "• 所有项目可用",
            "",
            "npx 方式",
            "• npx easy-opencode install",
            "• 无需全局安装",
            "• 临时使用场景"
        ],
        [
            "在项目中安装",
            "• eoc-install",
            "• 选择安装位置",
            "• Project-level: .opencode/",
            "• Global: 全局配置",
            "",
            "推荐：Project-level",
            "• 便于项目管理",
            "• 项目级配置",
            "• 独立依赖"
        ]
    )
    
    # 7. API Key 配置
    add_two_column_slide(
        prs,
        "第四章：DeepSeek API Key 配置",
        [
            "配置方式",
            "• 环境变量（推荐）",
            "  export DEEPSEEK_API_KEY",
            "• OpenCode 配置文件",
            "  ~/.config/opencode/...",
            "",
            "命令行临时配置",
            "• DEEPSEEK_API_KEY=...",
            "• 单次有效",
            "• 测试场景"
        ],
        [
            "安全注意事项 ⚠️",
            "• ❌ 永不提交到 Git",
            "• ✅ .env 添加到 .gitignore",
            "• ✅ 定期轮换 API Key",
            "• ✅ 不同项目用不同 Key",
            "",
            "验证配置",
            "• echo $DEEPSEEK_API_KEY",
            "• /test-connection",
            "• 检查连接状态"
        ]
    )
    
    # 8. EOC 命令
    add_command_slide(
        prs,
        "第五章：常用 EOC 命令",
        [
            "/plan - 创建实现计划",
            "/tdd - TDD 开发流程",
            "/code-review - 代码审查",
            "/security - 安全审查",
            "/build-fix - 修复构建错误",
            "/architect - 架构设计",
            "/e2e - 生成 E2E 测试",
            "/refactor-clean - 清理死代码"
        ]
    )
    
    # 9. 实战演示
    add_flow_slide(
        prs,
        "实战演示：30 分钟开发待办事项应用",
        [
            "规划阶段\n/plan\n(5分钟)",
            "TDD 开发\n/tdd\n(15分钟)",
            "安全审查\n/security\n(3分钟)",
            "代码审查\n/code-review\n(3分钟)",
            "E2E 测试\n/e2e\n(5分钟)"
        ]
    )
    
    # 10. 常见问题
    add_feature_boxes_slide(
        prs,
        "第六章：常见问题与解决方案",
        [
            "Node.js 安装失败 → 检查 PATH",
            "OpenCode 安装失败 → 检查网络",
            "API Key 不生效 → 验证环境变量",
            "Agent 不可用 → 重新安装 EOC",
            "命令行错误 → 检查拼写和参数",
            "连接失败 → 检查网络和 API Key"
        ]
    )
    
    # 11. 总结
    add_feature_boxes_slide(
        prs,
        "恭喜！您已掌握：",
        [
            "✅ Node.js 安装配置",
            "✅ OpenCode 安装验证",
            "✅ EOC 部署配置",
            "✅ API Key 安全设置",
            "✅ EOC 命令熟练使用",
            "✅ 工作流程完全掌握"
        ]
    )
    
    # 12. 下一步
    add_card_grid_slide(
        prs,
        "下一步行动",
        [
            "📚 立即实践",
            "🎯 深入学习",
            "💬 团队协作",
            "🔄 持续改进",
            "📊 追踪效果",
            "🚀 拓展应用"
        ]
    )
    
    # 13. Q&A
    add_qa_slide(
        prs,
        "Q & A",
        "感谢聆听！如有疑问，欢迎提问 💡"
    )
    
    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "opencode-training-presentation/OpenCode安装与配置指南-Coral版.pptx"
    prs.save(output_path)
    print(f"✅ 现代化演示文稿已生成: {output_path}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    print(f"🎨 设计风格: Coral Energy (Primary: Coral #F96167)")
    print(f"✨ 特性: 渐变背景、卡片布局、现代化排版")
