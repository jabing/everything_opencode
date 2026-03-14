#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
sys.path.insert(0, '/home/jabing/.local/lib/python3.14/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

# 创新大胆风格（Creative Bold - Google/Airbnb Style）
PRIMARY_COLOR = RGBColor(249, 97, 103)      # Coral
SECONDARY_COLOR = RGBColor(249, 231, 149)   # Gold
BACKGROUND_COLOR = RGBColor(47, 60, 126)    # Navy
TEXT_COLOR = RGBColor(255, 255, 255)        # White
HIGHLIGHT_COLOR = RGBColor(255, 255, 255)    # White cards

# 字体
TITLE_FONT = "Arial Black"
BODY_FONT = "Arial"

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title_slide(prs, title, subtitle, presenter):
    """添加现代化标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, PRIMARY_COLOR, BACKGROUND_COLOR)
    
    # 装饰圆形
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(2), Inches(2)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = SECONDARY_COLOR
    circle1.fill.fore_color.brightness = 0.3
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.5), Inches(3.5), Inches(2.5), Inches(2.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = SECONDARY_COLOR
    circle2.fill.fore_color.brightness = 0.2
    circle2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(84)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.font.name = BODY_FONT
    
    # 演讲者标签
    tag_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5), Inches(3), Inches(0.5)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = SECONDARY_COLOR
    tag_box.line.fill.background()
    
    tag_text = tag_box.text_frame
    tag_text.text = presenter
    tag_para = tag_text.paragraphs[0]
    tag_para.alignment = PP_ALIGN.CENTER
    tag_para.font.size = Pt(24)
    tag_para.font.bold = True
    tag_para.font.color.rgb = BACKGROUND_COLOR
    tag_para.font.name = BODY_FONT

def add_agenda_slide(prs, title, items):
    """添加目录页 - 卡片网格布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 2x3 卡片网格
    card_width = 4.5
    card_height = 1.3
    positions = [
        (Inches(0.5), Inches(1.5)),
        (Inches(5), Inches(1.5)),
        (Inches(0.5), Inches(3)),
        (Inches(5), Inches(3)),
        (Inches(0.5), Inches(4.5)),
        (Inches(5), Inches(4.5))
    ]
    
    for i, (item, pos) in enumerate(zip(items, positions)):
        x, y = pos
        
        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = HIGHLIGHT_COLOR
        card.line.color.rgb = PRIMARY_COLOR
        card.line.width = Pt(3)
        
        # 编号圆圈
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.35), Inches(0.6), Inches(0.6)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = PRIMARY_COLOR
        num_circle.line.fill.background()
        
        # 编号
        num_text = num_circle.text_frame
        num_text.text = str(i + 1)
        num_para = num_text.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(24)
        num_para.font.bold = True
        num_para.font.color.rgb = TEXT_COLOR
        num_para.font.name = BODY_FONT
        
        # 文本
        text_frame = card.text_frame
        text_frame.text = item
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(24)
        text_para.font.color.rgb = BACKGROUND_COLOR
        text_para.font.name = BODY_FONT

def add_two_column_slide(prs, title, left_content, right_content):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, RGBColor(240, 240, 240), RGBColor(255, 255, 255))
    
    # 标题栏
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.2), Inches(9.6), Inches(0.8)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = PRIMARY_COLOR
    title_box.line.fill.background()
    
    title_text = title_box.text_frame
    title_text.text = title
    title_para = title_text.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.font.name = TITLE_FONT
    
    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.4), Inches(3.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, line in enumerate(left_content):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = BACKGROUND_COLOR
        p.font.name = BODY_FONT
        p.space_before = Pt(8)
        p.space_after = Pt(6)
    
    # 分隔线
    divider = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.5), Inches(0.1), Inches(3.8)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = SECONDARY_COLOR
    divider.fill.fore_color.brightness = 0.5
    divider.line.fill.background()
    
    # 右栏
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(3.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, line in enumerate(right_content):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = BACKGROUND_COLOR
        p.font.name = BODY_FONT
        p.space_before = Pt(8)
        p.space_after = Pt(6)

def add_feature_cards_slide(prs, title, features):
    """添加特性卡片页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(52)
    title_para.font.bold = True
    title_para.font.color.rgb = SECONDARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 特性卡片
    for i, feature in enumerate(features):
        row = i // 2
        col = i % 2
        
        x = Inches(0.5) if col == 0 else Inches(5)
        y = Inches(1.3) + Inches(row * 1.4)
        
        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = HIGHLIGHT_COLOR
        card.line.color.rgb = SECONDARY_COLOR
        card.line.width = Pt(2)
        
        # 特性图标（用圆圈代替）
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.35), Inches(0.5), Inches(0.5)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = PRIMARY_COLOR
        icon.line.fill.background()
        
        # 文本
        text_frame = card.text_frame
        text_frame.text = f"  {feature}"
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(24)
        text_para.font.bold = True
        text_para.font.color.rgb = BACKGROUND_COLOR
        text_para.font.name = BODY_FONT

def add_command_bars_slide(prs, title, commands):
    """添加命令条页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, RGBColor(47, 60, 126), RGBColor(30, 40, 90))
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = SECONDARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 命令条
    y_pos = 1.2
    for i, command in enumerate(commands):
        # 命令条背景
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.7)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = HIGHLIGHT_COLOR
        bar.line.color.rgb = SECONDARY_COLOR
        bar.line.width = Pt(2)
        
        # 编号
        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), Inches(y_pos + 0.15), Inches(0.4), Inches(0.4)
        )
        num.fill.solid()
        num.fill.fore_color.rgb = PRIMARY_COLOR
        num.line.fill.background()
        
        num_text = num.text_frame
        num_text.text = str(i + 1)
        num_para = num_text.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(18)
        num_para.font.bold = True
        num_para.font.color.rgb = TEXT_COLOR
        num_para.font.name = BODY_FONT
        
        # 命令文本
        text_frame = bar.text_frame
        text_frame.text = f"  {command}"
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(24)
        text_para.font.bold = True
        text_para.font.color.rgb = BACKGROUND_COLOR
        text_para.font.name = BODY_FONT
        
        y_pos += 0.85

def add_timeline_slide(prs, title, steps):
    """添加时间轴/流程页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 时间轴节点
    x_pos = 0.5
    for i, step in enumerate(steps):
        # 节点圆形
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x_pos), Inches(1.8), Inches(1), Inches(1)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = PRIMARY_COLOR
        node.line.color.rgb = SECONDARY_COLOR
        node.line.width = Pt(3)
        
        # 节点编号
        num_text = node.text_frame
        num_text.text = str(i + 1)
        num_para = num_text.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(36)
        num_para.font.bold = True
        num_para.font.color.rgb = TEXT_COLOR
        num_para.font.name = BODY_FONT
        
        # 步骤描述
        desc_box = slide.shapes.add_textbox(Inches(x_pos - 0.2), Inches(3), Inches(1.4), Inches(1.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = step
        desc_para = desc_frame.paragraphs[0]
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = BACKGROUND_COLOR
        desc_para.font.name = BODY_FONT
        
        x_pos += 1.9
        
        # 连接线（除了最后一个）
        if i < len(steps) - 1:
            line = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(x_pos - 0.4), Inches(2.3),
                Inches(x_pos - 0.4 + 1.5), Inches(2.3)
            )
            line.line.color.rgb = SECONDARY_COLOR
            line.line.width = Pt(4)

def add_icon_boxes_slide(prs, title, items):
    """添加图标盒子页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, PRIMARY_COLOR, BACKGROUND_COLOR)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(52)
    title_para.font.bold = True
    title_para.font.color.rgb = SECONDARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 3x2 图标盒子
    positions = [
        (Inches(0.5), Inches(1.5)),
        (Inches(3.5), Inches(1.5)),
        (Inches(6.5), Inches(1.5)),
        (Inches(0.5), Inches(3.2)),
        (Inches(3.5), Inches(3.2)),
        (Inches(6.5), Inches(3.2))
    ]
    
    for item, pos in zip(items, positions):
        x, y = pos
        
        # 盒子
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(1.4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = HIGHLIGHT_COLOR
        box.line.color.rgb = SECONDARY_COLOR
        box.line.width = Pt(2)
        
        # 图标圆圈
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1), y + Inches(0.2), Inches(0.8), Inches(0.8)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = PRIMARY_COLOR
        icon.line.fill.background()
        
        text_frame.text = f"\n{item}"
        text_para = text_frame.paragraphs[0]
        text_para.alignment = PP_ALIGN.CENTER
        text_para.font.size = Pt(18)
        text_para.font.bold = True
        text_para.font.color.rgb = BACKGROUND_COLOR
        text_para.font.name = BODY_FONT

def add_qa_slide(prs, title, subtitle):
    """添加问答页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变背景
    add_gradient_background(slide, PRIMARY_COLOR, RGBColor(47, 60, 126))
    
    # 装饰元素
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.5), Inches(3), Inches(2.5), Inches(2.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY_COLOR
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(88)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.font.name = BODY_FONT

def create_modern_presentation():
    """创建现代化演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 标题页
    add_title_slide(
        prs,
        "OpenCode 安装\n与配置指南",
        "从零开始打造 AI 编码助手",
        "🎯 部门内部培训"
    )
    
    # 2. 目录
    add_agenda_slide(
        prs,
        "课程大纲",
        [
            "Node.js 安装与配置",
            "OpenCode 安装",
            "Easy-OpenCode 部署",
            "DeepSeek API Key 配置",
            "开始使用 Easy-OpenCode",
            "常见问题与解决方案"
        ]
    )
    
    # 3. 第一章：Node.js
    add_two_column_slide(
        prs,
        "第一章：Node.js 安装",
        [
            "什么是 Node.js？",
            "• Chrome V8 引擎运行时",
            "• 跨平台支持",
            "• npm 包管理器",
            "• OpenCode 基础",
            "",
            "官方安装（推荐）",
            "• 访问 nodejs.org",
            "• 下载 LTS 版本",
            "• 运行安装程序"
        ],
        [
            "Linux 命令行安装",
            "• Ubuntu/Debian:",
            "  curl -fsSL https://deb.nodesource.com/setup_lts.x | sudo -E bash -",
            "• sudo apt-get install -y nodejs",
            "",
            "验证安装",
            "• node --version",
            "• npm --version"
        ]
    )
    
    # 4. 第二章：OpenCode
    add_two_column_slide(
        prs,
        "第二章：OpenCode 安装",
        [
            "OpenCode 简介",
            "• 强大的 AI 编码助手",
            "• 13 个专业 AI Agent",
            "• 50+ 工作流技能",
            "• 33 个命令",
            "",
            "Homebrew 安装",
            "• brew install opencode",
            "• 适用于 macOS/Linux",
            "• 自动配置环境变量"
        ],
        [
            "NPM 全局安装",
            "• npm install -g opencode",
            "• 适用于所有平台",
            "• 验证安装",
            "• opencode --version",
            "",
            "验证成功标志",
            "• 显示帮助信息",
            "• Agent 列表可见",
            "• 命令列表完整"
        ]
    )
    
    # 5. 第三章：Easy-OpenCode 介绍
    add_feature_cards_slide(
        prs,
        "第三章：Easy-OpenCode 是什么？",
        [
            "🤖 14 个专业 AI Agent",
            "⚡ 50+ 工作流技能",
            "🛠️ 33 个快捷命令",
            "🔗 自动化钩子工作流",
            "📦 一键安装开箱即用",
            "🎯 TDD 驱动确保质量"
        ]
    )
    
    # 6. EOC 安装
    add_two_column_slide(
        prs,
        "Easy-OpenCode 安装方式",
        [
            "全局安装（推荐）⭐",
            "• npm install -g easy-opencode",
            "• eoc-install --version",
            "• 所有项目可用",
            "",
            "npx 方式",
            "• npx easy-opencode install",
            "• 无需全局安装",
            "• 临时使用场景"
        ],
        [
            "在项目中安装",
            "• eoc-install",
            "• 选择安装位置",
            "• Project-level: .opencode/",
            "• Global: 全局配置目录",
            "",
            "推荐：Project-level",
            "• 便于项目管理",
            "• 项目级配置",
            "• 独立依赖"
        ]
    )
    
    # 7. API Key 配置
    add_two_column_slide(
        prs,
        "第四章：DeepSeek API Key 配置",
        [
            "配置方式",
            "• 环境变量（推荐）",
            "  export DEEPSEEK_API_KEY=\"key\"",
            "• OpenCode 配置文件",
            "  ~/.config/opencode/opencode.json",
            "",
            "命令行临时配置",
            "• DEEPSEEK_API_KEY=\"key\" opencode",
            "• 单次有效",
            "• 测试场景"
        ],
        [
            "安全注意事项 ⚠️",
            "• ❌ 永不提交到 Git",
            "• ✅ .env 添加到 .gitignore",
            "• ✅ 定期轮换 API Key",
            "• ✅ 不同项目用不同 Key",
            "",
            "验证配置",
            "• echo $DEEPSEEK_API_KEY",
            "• /test-connection",
            "• 检查连接状态"
        ]
    )
    
    # 8. EOC 命令
    add_command_bars_slide(
        prs,
        "第五章：常用 EOC 命令",
        [
            "/plan - 创建实现计划",
            "/tdd - TDD 开发流程",
            "/code-review - 代码审查",
            "/security - 安全审查",
            "/build-fix - 修复构建错误",
            "/architect - 架构设计",
            "/e2e - 生成 E2E 测试",
            "/refactor-clean - 清理死代码"
        ]
    )
    
    # 9. 实战演示
    add_timeline_slide(
        prs,
        "实战演示：30 分钟开发待办事项应用",
        [
            "规划阶段\n/plan\n(5分钟)",
            "TDD 开发\n/tdd\n(15分钟)",
            "安全审查\n/security\n(3分钟)",
            "代码审查\n/code-review\n(3分钟)",
            "E2E 测试\n/e2e\n(5分钟)"
        ]
    )
    
    # 10. 常见问题
    add_icon_boxes_slide(
        prs,
        "第六章：常见问题与解决方案",
        [
            "Node.js\n安装失败",
            "检查 PATH\n环境变量",
            "OpenCode\n安装失败",
            "检查网络\nsudo 安装",
            "API Key\n不生效",
            "验证环境\n重启 OpenCode"
        ]
    )
    
    # 11. 总结
    add_icon_boxes_slide(
        prs,
        "恭喜！您已掌握：",
        [
            "✅ Node.js\n安装配置",
            "✅ OpenCode\n安装验证",
            "✅ EOC\n部署配置",
            "✅ API Key\n安全设置",
            "✅ EOC 命令\n熟练使用",
            "✅ 工作流程\n完全掌握"
        ]
    )
    
    # 12. 下一步
    add_feature_cards_slide(
        prs,
        "下一步行动",
        [
            "📚 立即实践\n创建测试项目",
            "🎯 深入学习\n阅读 AGENTS.md",
            "💬 团队协作\n分享建立实践",
            "🔄 持续改进\n定期技术交流",
            "📊 追踪效果\n测量效率提升",
            "🚀 拓展应用\n探索更多技能"
        ]
    )
    
    # 13. Q&A
    add_qa_slide(
        prs,
        "Q & A",
        "感谢聆听！如有疑问，欢迎提问 💡"
    )
    
    return prs

if __name__ == "__main__":
    prs = create_modern_presentation()
    output_path = "opencode-training-presentation/OpenCode安装与配置指南-现代版.pptx"
    prs.save(output_path)
    print(f"✅ 现代化演示文稿已生成: {output_path}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    print(f"🎨 设计风格: Creative Bold (Coral + Gold)")
    print(f"📏 尺寸: 16:9 ({prs.slide_width} x {prs.slide_height})")
