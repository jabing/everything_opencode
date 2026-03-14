#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
sys.path.insert(0, '/home/jabing/.local/lib/python3.14/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 颜色定义 - Midnight Executive 主题
PRIMARY_COLOR = RGBColor(30, 39, 97)    # Navy
SECONDARY_COLOR = RGBColor(202, 220, 252)  # Ice blue
ACCENT_COLOR = RGBColor(255, 215, 0)     # Gold
BACKGROUND_COLOR = RGBColor(255, 255, 255)  # White

# 字体定义
TITLE_FONT = "Georgia"
BODY_FONT = "Calibri"

def add_title_slide(prs, title, subtitle, presenter):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.font.name = BODY_FONT
    
    # 演讲者
    presenter_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = presenter
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.alignment = PP_ALIGN.CENTER
    presenter_para.font.size = Pt(24)
    presenter_para.font.color.rgb = RGBColor(200, 200, 200)
    presenter_para.font.name = BODY_FONT

def add_content_slide(prs, title, bullets):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = PRIMARY_COLOR
    title_box.line.fill.background()
    
    title_text = title_box.text_frame
    title_text.text = title
    title_text.word_wrap = True
    title_para = title_text.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = TITLE_FONT
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]
        
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(30, 39, 97)
        p.font.name = BODY_FONT
        p.space_before = Pt(12)
        p.space_after = Pt(6)

def add_highlight_slide(prs, title, items):
    """添加高亮页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 内容卡片
    y_pos = 1.8
    for item in items:
        # 卡片背景
        card_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y_pos), Inches(8), Inches(1)
        )
        card_box.fill.solid()
        card_box.fill.fore_color.rgb = SECONDARY_COLOR
        card_box.line.color.rgb = PRIMARY_COLOR
        card_box.line.width = Pt(2)
        
        # 文本
        text_frame = card_box.text_frame
        text_frame.text = item
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.font.name = BODY_FONT
        
        y_pos += 1.2

def add_demo_slide(prs, title, steps, summary):
    """添加演示页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.name = TITLE_FONT
    
    # 步骤
    y_pos = 1.3
    for step_info in steps:
        step = step_info["step"]
        command = step_info["command"]
        output = step_info["output"]
        
        # 步骤框
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(1)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = ACCENT_COLOR
        step_box.line.color.rgb = PRIMARY_COLOR
        
        step_text = step_box.text_frame
        step_text.text = step
        step_para = step_text.paragraphs[0]
        step_para.alignment = PP_ALIGN.CENTER
        step_para.font.size = Pt(24)
        step_para.font.bold = True
        step_para.font.color.rgb = PRIMARY_COLOR
        step_para.font.name = BODY_FONT
        
        # 命令
        cmd_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.7), Inches(8), Inches(0.4))
        cmd_frame = cmd_box.text_frame
        cmd_frame.text = f"命令: {command}"
        cmd_para = cmd_frame.paragraphs[0]
        cmd_para.font.size = Pt(16)
        cmd_para.font.italic = True
        cmd_para.font.color.rgb = RGBColor(100, 100, 100)
        cmd_para.font.name = "Courier New"
        
        # 输出
        out_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 1.1), Inches(8), Inches(0.4))
        out_frame = out_box.text_frame
        out_frame.text = f"输出: {output}"
        out_para = out_frame.paragraphs[0]
        out_para.font.size = Pt(16)
        out_para.font.color.rgb = RGBColor(30, 39, 97)
        out_para.font.name = BODY_FONT
        
        y_pos += 1.6
    
    # 总结
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y_pos + 0.2), Inches(8), Inches(0.6)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = PRIMARY_COLOR
    summary_box.line.color.rgb = ACCENT_COLOR
    
    summary_text = summary_box.text_frame
    summary_text.text = summary
    summary_para = summary_text.paragraphs[0]
    summary_para.alignment = PP_ALIGN.CENTER
    summary_para.font.size = Pt(28)
    summary_para.font.bold = True
    summary_para.font.color.rgb = RGBColor(255, 255, 255)
    summary_para.font.name = BODY_FONT

def add_qa_slide(prs, title, subtitle):
    """添加问答页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = TITLE_FONT
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = SECONDARY_COLOR
    subtitle_para.font.name = BODY_FONT

def create_presentation():
    """创建完整的演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 比例
    
    # 1. 标题页
    add_title_slide(
        prs,
        "OpenCode 安装与配置指南",
        "从零开始打造 AI 编码助手",
        "部门内部培训"
    )
    
    # 2. 目录
    add_content_slide(
        prs,
        "课程大纲",
        [
            "1. Node.js 安装与配置",
            "2. OpenCode 安装",
            "3. Easy-OpenCode 部署",
            "4. DeepSeek API Key 配置",
            "5. 开始使用 Easy-OpenCode",
            "6. 常见问题与解决方案"
        ]
    )
    
    # 3. 第一章：Node.js
    add_content_slide(
        prs,
        "第一章：Node.js 安装",
        [
            "什么是 Node.js？",
            "• 基于 Chrome V8 引擎的 JavaScript 运行时",
            "• 跨平台支持（Windows/macOS/Linux）",
            "• 强大的包管理器（npm）",
            "• OpenCode 的运行基础",
            "",
            "安装方式：",
            "• 方式一：官方安装包（推荐）",
            "  访问 https://nodejs.org/，下载 LTS 版本",
            "• 方式二：命令行安装（Linux）",
            "  Ubuntu: curl -fsSL https://deb.nodesource.com/setup_lts.x | sudo -E bash -",
            "• 验证: node --version, npm --version"
        ]
    )
    
    # 4. 第二章：OpenCode
    add_content_slide(
        prs,
        "第二章：OpenCode 安装",
        [
            "OpenCode 简介",
            "• 强大的 AI 编码助手",
            "• 13 个专业 AI Agent",
            "• 50+ 工作流技能",
            "• 33 个命令",
            "• 自动化工作流",
            "",
            "安装 OpenCode",
            "• 方式一：Homebrew（macOS/Linux）",
            "  brew install opencode",
            "• 方式二：NPM（所有平台）",
            "  npm install -g opencode",
            "• 验证安装: opencode --version"
        ]
    )
    
    # 5. 第三章：Easy-OpenCode
    add_content_slide(
        prs,
        "第三章：Easy-OpenCode 部署",
        [
            "Easy-OpenCode 是什么？",
            "• 面向 OpenCode 的生产级 AI 编码插件",
            "• 14 个专业 AI Agent",
            "• 50+ 工作流技能",
            "• 33 个命令",
            "• 自动化钩子工作流",
            "",
            "安装 Easy-OpenCode",
            "• 方式一：全局安装（推荐）⭐",
            "  npm install -g easy-opencode",
            "  eoc-install --version",
            "• 方式二：使用 npx",
            "  npx easy-opencode install",
            "• 在项目中安装: eoc-install"
        ]
    )
    
    # 6. 核心优势
    add_highlight_slide(
        prs,
        "核心优势",
        [
            "📦 一键安装，开箱即用",
            "🎯 TDD 驱动开发，确保代码质量",
            "🔒 内置安全审查，防止漏洞",
            "🚀 自动化工作流，提升开发效率"
        ]
    )
    
    # 7. 第四章：DeepSeek API Key
    add_content_slide(
        prs,
        "第四章：DeepSeek API Key 配置",
        [
            "配置方式",
            "• 方式一：环境变量（推荐）",
            "  export DEEPSEEK_API_KEY=\"your-api-key-here\"",
            "  source ~/.bashrc",
            "• 方式二：OpenCode 配置文件（持久化）",
            "  编辑 ~/.config/opencode/opencode.json",
            "• 方式三：命令行临时配置",
            "  DEEPSEEK_API_KEY=\"key\" opencode",
            "",
            "安全注意事项 ⚠️",
            "• 永远不要将 API Key 提交到 Git",
            "• 将 .env 添加到 .gitignore",
            "• 定期轮换 API Key"
        ]
    )
    
    # 8. 第五章：常用命令
    add_content_slide(
        prs,
        "第五章：常用 EOC 命令",
        [
            "/plan - 创建实现计划",
            "/tdd - TDD 开发流程",
            "/code-review - 代码审查",
            "/security - 安全审查",
            "/build-fix - 修复构建错误",
            "/architect - 架构设计",
            "/e2e - 生成 E2E 测试",
            "/refactor-clean - 清理死代码"
        ]
    )
    
    # 9. 实战演示
    add_demo_slide(
        prs,
        "实战演示：开发待办事项应用",
        [
            {
                "step": "1️⃣ 规划阶段",
                "command": "/plan 创建一个待办事项应用",
                "output": "实现计划已生成（5分钟）"
            },
            {
                "step": "2️⃣ TDD 开发",
                "command": "/tdd 实现 TodoList 组件",
                "output": "测试+实现+重构（15分钟）"
            },
            {
                "step": "3️⃣ 安全审查",
                "command": "/security 检查安全性",
                "output": "安全检查报告（3分钟）"
            },
            {
                "step": "4️⃣ 代码审查",
                "command": "/code-review 审查代码质量",
                "output": "代码质量报告（3分钟）"
            }
        ],
        "总计：约 30 分钟完成从 0 到 1 的开发！🚀"
    )
    
    # 10. 第六章：常见问题
    add_content_slide(
        prs,
        "第六章：常见问题与解决方案",
        [
            "Q1: Node.js 安装失败",
            "  A: 检查 PATH 环境变量，重新运行安装程序",
            "",
            "Q2: OpenCode 安装失败",
            "  A: 检查网络连接，使用 sudo 安装，清除 npm 缓存",
            "",
            "Q3: API Key 不生效",
            "  A: 验证环境变量，检查 .env 文件，重启 OpenCode",
            "",
            "Q4: Agent 不可用",
            "  A: 确认 EOC 安装成功，重新安装，清除缓存"
        ]
    )
    
    # 11. 总结
    add_highlight_slide(
        prs,
        "恭喜！您已掌握：",
        [
            "✅ Node.js 安装与配置",
            "✅ OpenCode 安装与验证",
            "✅ Easy-OpenCode 部署",
            "✅ DeepSeek API Key 配置",
            "✅ EOC 命令和操作掌握"
        ]
    )
    
    # 12. 下一步
    add_highlight_slide(
        prs,
        "下一步行动",
        [
            "📚 立即实践：创建测试项目，尝试第一个命令",
            "🎯 深入学习：阅读 AGENTS.md，探索 skills/ 目录",
            "💬 团队协作：分享经验，建立最佳实践"
        ]
    )
    
    # 13. Q&A
    add_qa_slide(
        prs,
        "Q & A",
        "感谢聆听！如有疑问，欢迎提问"
    )
    
    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "opencode-training-presentation/OpenCode安装与配置指南.pptx"
    prs.save(output_path)
    print(f"✅ 演示文稿已生成: {output_path}")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
