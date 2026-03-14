#!/usr/bin/env python3
import sys
sys.path.insert(0, '/home/jabing/.local/lib/python3.14/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(10, 10, 10)
CARD = RGBColor(30, 30, 35)
PRI = RGBColor(0, 212, 255)
ACC = RGBColor(0, 255, 255)
HIL = RGBColor(255, 0, 255)
TXT = RGBColor(255, 255, 255)
CODE = RGBColor(0, 255, 136)

TF = "Arial Black"
BF = "Arial"
CF = "Consolas"

def db(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG

def ts(prs, title, subtitle, tag_txt):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    for i in range(5):
        line = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3 + i * 0.15), Inches(0.5 + i * 0.3), Inches(9.4 - i * 0.1), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = ACC
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.8))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb.text_frame.paragraphs[0].font.size = Pt(76)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PRI
    tb.text_frame.paragraphs[0].font.name = TF
    sb = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
    sb.text_frame.text = subtitle
    sb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sb.text_frame.paragraphs[0].font.size = Pt(28)
    sb.text_frame.paragraphs[0].font.color.rgb = HIL
    sb.text_frame.paragraphs[0].font.name = BF
    tag_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.5), Inches(4), Inches(0.6))
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = CARD
    tag_shape.line.color.rgb = ACC
    tag_shape.line.width = Pt(2)
    tag_shape.text_frame.text = tag_txt
    tag_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_shape.text_frame.paragraphs[0].font.size = Pt(24)
    tag_shape.text_frame.paragraphs[0].font.bold = True
    tag_shape.text_frame.paragraphs[0].font.color.rgb = ACC
    tag_shape.text_frame.paragraphs[0].font.name = TF

def cs(prs, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(4.5), Inches(1.5), Inches(1.5))
    c.fill.solid()
    c.fill.fore_color.rgb = PRI
    c.fill.fore_color.brightness = 0.2
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CARD
    bar.line.color.rgb = PRI
    bar.line.width = Pt(2)
    bar.text_frame.text = title
    bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    bar.text_frame.paragraphs[0].font.size = Pt(40)
    bar.text_frame.paragraphs[0].font.bold = True
    bar.text_frame.paragraphs[0].font.color.rgb = PRI
    bar.text_frame.paragraphs[0].font.name = TF
    y = 1.3
    for item in items:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = ACC
        card.line.width = Pt(1)
        card.text_frame.text = item
        card.text_frame.paragraphs[0].font.size = Pt(22)
        card.text_frame.paragraphs[0].font.color.rgb = TXT
        card.text_frame.paragraphs[0].font.name = BF
        y += 0.9

def cds(prs, title, blocks):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(44)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PRI
    tb.text_frame.paragraphs[0].font.name = TF
    y = 1.2
    for block in blocks:
        cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.6))
        cb.fill.solid()
        cb.fill.fore_color.rgb = RGBColor(20, 20, 20)
        cb.line.color.rgb = CODE
        cb.line.width = Pt(2)
        cb.text_frame.text = block
        cb.text_frame.paragraphs[0].font.size = Pt(16)
        cb.text_frame.paragraphs[0].font.color.rgb = CODE
        cb.text_frame.paragraphs[0].font.name = CF
        y += 0.75

def fs(prs, title, features):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(48)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PRI
    tb.text_frame.paragraphs[0].font.name = TF
    pos = [(Inches(0.5), Inches(1.2)), (Inches(5), Inches(1.2)), (Inches(0.5), Inches(2.4)), (Inches(5), Inches(2.4)), (Inches(0.5), Inches(3.6)), (Inches(5), Inches(3.6))]
    for f, p in zip(features, pos):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p[0], p[1], Inches(4.5), Inches(1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = ACC
        card.line.width = Pt(1)
        card.text_frame.text = f"  {f}"
        card.text_frame.paragraphs[0].font.size = Pt(24)
        card.text_frame.paragraphs[0].font.bold = True
        card.text_frame.paragraphs[0].font.color.rgb = TXT
        card.text_frame.paragraphs[0].font.name = BF

def fls(prs, title, steps):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(40)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PRI
    tb.text_frame.paragraphs[0].font.name = TF
    x = 0.5
    for i, step_txt in enumerate(steps):
        node = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(1.8), Inches(2.2))
        node.fill.solid()
        node.fill.fore_color.rgb = CARD
        node.line.color.rgb = HIL
        node.line.width = Pt(2)
        node.text_frame.word_wrap = True
        node.text_frame.text = step_txt
        node.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        node.text_frame.paragraphs[0].font.size = Pt(16)
        node.text_frame.paragraphs[0].font.bold = True
        node.text_frame.paragraphs[0].font.color.rgb = TXT
        node.text_frame.paragraphs[0].font.name = BF
        x += 1.9
        if i < len(steps) - 1:
            conn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.3), Inches(2.5), Inches(0.5), Inches(0.1))
            conn.fill.solid()
            conn.fill.fore_color.rgb = ACC

def rs(prs, title, resources):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(48)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PRI
    tb.text_frame.paragraphs[0].font.name = TF
    y = 1.3
    for name, url in resources:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = HIL
        card.line.width = Pt(2)
        card.text_frame.text = name
        card.text_frame.paragraphs[0].font.size = Pt(20)
        card.text_frame.paragraphs[0].font.bold = True
        card.text_frame.paragraphs[0].font.color.rgb = PRI
        card.text_frame.paragraphs[0].font.name = BF
        up = card.text_frame.add_paragraph()
        up.text = url
        up.font.size = Pt(14)
        up.font.color.rgb = CODE
        up.font.italic = True
        up.font.name = CF
        y += 1.0

def qs(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    db(s)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(2.5), Inches(2.5), Inches(2.5))
    c.fill.solid()
    c.fill.fore_color.rgb = HIL
    c.fill.fore_color.brightness = 0.3
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb.text_frame.paragraphs[0].font.size = Pt(88)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PRI
    tb.text_frame.paragraphs[0].font.name = TF
    sb = s.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    sb.text_frame.text = subtitle
    sb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sb.text_frame.paragraphs[0].font.size = Pt(32)
    sb.text_frame.paragraphs[0].font.color.rgb = ACC
    sb.text_frame.paragraphs[0].font.name = BF

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

ts(prs, "OpenCode 安装\n与配置指南", "从零开始打造 AI 编码助手", "< 部门内部培训 >")

cs(prs, "课程大纲", [
    "1. Node.js 安装",
    "2. OpenCode 安装",
    "3. Easy-OpenCode 部署",
    "4. DeepSeek API Key 配置",
    "5. 开始使用 Easy-OpenCode",
    "6. 常见问题与解决方案",
    "7. 相关资源与文档"
])

cds(prs, "第一章：Node.js 安装", [
    "# 官方安装包（推荐）",
    "访问: https://nodejs.org/",
    "下载 LTS 版本，选择对应系统的安装包",
    "Windows: 双击 .msi 文件",
    "macOS: 双击 .pkg 文件",
    "Linux: 按照官网指南操作",
    "",
    "# 验证安装",
    "node --version",
    "npm --version"
])

cds(prs, "第二章：OpenCode 安装", [
    "# 全局安装（推荐）",
    "npm install -g opencode",
    "",
    "# 验证安装",
    "opencode --version",
    "opencode --help"
])

fs(prs, "第三章：Easy-OpenCode 核心特性", [
    "14 个专业 AI Agent",
    "50+ 工作流技能",
    "33 个快捷命令",
    "自动化钩子工作流",
    "TDD 驱动开发",
    "内置安全审查"
])

cds(prs, "Easy-OpenCode 安装", [
    "# 全局安装（推荐）",
    "npm install -g easy-opencode",
    "",
    "# 在项目中安装",
    "eoc-install",
    "或",
    "npx easy-opencode install"
])

cds(prs, "第四章：DeepSeek API Key 配置", [
    "# 环境变量配置（推荐）",
    "export DEEPSEEK_API_KEY=\"your-api-key-here\"",
    "source ~/.bashrc  或 source ~/.zshrc",
    "",
    "# OpenCode 配置文件（持久化）",
    "# 编辑: ~/.config/opencode/opencode.json",
    "# 添加:",
    "{\"env\": {\"DEEPSEEK_API_KEY\": \"your-api-key-here\"}}",
    "",
    "# 安全注意事项",
    "# 将 .env 添加到 .gitignore",
    "# 定期轮换 API Key"
])

cs(prs, "第五章：常用 EOC 命令", [
    "/plan - 创建实现计划",
    "/tdd - TDD 开发流程",
    "/code-review - 代码审查",
    "/security - 安全审查",
    "/build-fix - 修复构建错误",
    "/architect - 架构设计",
    "/e2e - 生成 E2E 测试",
    "/refactor-clean - 清理死代码"
])

fls(prs, "实战演示：30 分钟完整开发流程", [
    "步骤 1: 规划\n/plan\n5分钟",
    "步骤 2: TDD 开发\n/tdd\n15分钟",
    "步骤 3: 安全审查\n/security\n3分钟",
    "步骤 4: 代码审查\n/code-review\n3分钟",
    "步骤 5: E2E 测试\n/e2e\n5分钟"
])

cs(prs, "第六章：常见问题与解决方案", [
    "Q: Node.js 安装失败",
    "A: 检查 PATH 环境变量，重新运行安装程序",
    "Q: OpenCode 安装失败",
    "A: 检查网络连接，清除 npm 缓存后重试",
    "Q: API Key 不生效",
    "A: 验证环境变量，重启 OpenCode",
    "Q: Agent 不可用",
    "A: 确认 EOC 安装成功，重新安装"
])

resources = [
    ("Node.js 官方网站", "https://nodejs.org/"),
    ("OpenCode GitHub 仓库", "https://github.com/opencode-ai/opencode"),
    ("Easy-OpenCode GitHub 仓库", "https://github.com/jabing/easy_opencode"),
    ("DeepSeek 平台", "https://platform.deepseek.com/")
]
rs(prs, "相关资源与文档", resources)

fs(prs, "恭喜！您已掌握：", [
    "Node.js 安装配置",
    "OpenCode 安装验证",
    "EOC 部署配置",
    "API Key 安全设置",
    "EOC 命令熟练使用",
    "工作流程完全掌握"
])

cs(prs, "下一步行动", [
    "🚀 立即实践：创建测试项目，尝试第一个命令",
    "📚 深入学习：阅读 AGENTS.md，探索 skills/ 目录",
    "💬 团队协作：分享经验，建立最佳实践",
    "🔄 持续改进：定期技术交流，追踪效率提升"
])

qs(prs, "Q & A", "感谢聆听！如有疑问，欢迎提问 💡")

prs.save("opencode-training-presentation/OpenCode安装与配置指南-最终版.pptx")
print("Done! 14 slides created with dark tech style and resource links.")
